"""Generate class presentation PPTX — revised v4."""
import re, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = os.path.dirname(os.path.abspath(__file__))
CKPT = os.path.join(BASE, "checkpoints")

BG      = RGBColor(0x0F, 0x17, 0x2A)
PANEL   = RGBColor(0x16, 0x24, 0x3E)
ACCENT  = RGBColor(0x4A, 0x9E, 0xFF)
ACCENT2 = RGBColor(0x56, 0xD3, 0xA0)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xCC, 0xD6, 0xE8)
YELLOW  = RGBColor(0xFF, 0xD7, 0x5E)
RED     = RGBColor(0xFF, 0x6B, 0x6B)
ORANGE  = RGBColor(0xFF, 0xA0, 0x40)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

_NUM = re.compile(r'([+\-]?\d[\d,\.]*(?:%|×|ns|ms)?)')

def lerp_rgb(c1, c2, t):
    return RGBColor(*(int(a + t*(b-a)) for a, b in zip(c1, c2)))

# ── Primitives ────────────────────────────────────────────────────────────────
def bg(slide, color=BG):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def rect(slide, l, t, w, h, fc, lc=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fc
    if lc: s.line.color.rgb = lc
    else:  s.line.fill.background()
    return s

def box(slide, l, t, w, h, text, size=16, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf.word_wrap = True
    p = tf.text_frame.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tf

def img(slide, path, l, t, w, h=None):
    if not os.path.exists(path): return
    kw = dict(left=Inches(l), top=Inches(t), width=Inches(w))
    if h: kw['height'] = Inches(h)
    slide.shapes.add_picture(path, **kw)

def hline(slide, y, color=ACCENT, thickness=Pt(1.5)):
    c = slide.shapes.add_connector(1, Inches(0.4), Inches(y), Inches(12.9), Inches(y))
    c.line.color.rgb = color; c.line.width = thickness

# ── Rich text ─────────────────────────────────────────────────────────────────
def _rich_runs(para, text, base, num_color=YELLOW, size=16, bold=False, italic=False):
    for part in _NUM.split(text):
        if not part: continue
        r = para.add_run(); r.text = part
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = num_color if _NUM.fullmatch(part) else base

def rich_box(slide, l, t, w, h, text, size=16, bold=False,
             base=WHITE, num_color=YELLOW, align=PP_ALIGN.LEFT):
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf.word_wrap = True
    p = tf.text_frame.paragraphs[0]; p.alignment = align
    _rich_runs(p, text, base=base, num_color=num_color, size=size, bold=bold)

def bullet_box(slide, l, t, w, h, items, size=16, color=LGRAY, rich=False, num_color=YELLOW):
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf.word_wrap = True; tf.text_frame.auto_size = None
    first = True
    for item in items:
        level, text = item if isinstance(item, tuple) else (0, item)
        para = tf.text_frame.paragraphs[0] if first else tf.text_frame.add_paragraph()
        first = False; para.level = level; para.space_after = Pt(3)
        bullet = "▸ " if level == 0 else "  • "
        if rich:
            r0 = para.add_run(); r0.text = bullet
            r0.font.size = Pt(size); r0.font.color.rgb = color
            _rich_runs(para, text, base=color, num_color=num_color, size=size)
        else:
            r = para.add_run(); r.text = bullet + text
            r.font.size = Pt(size); r.font.color.rgb = color

# ── Slide templates ───────────────────────────────────────────────────────────
def section_header(slide, title, subtitle=None):
    bg(slide); rect(slide, 0, 0, 0.07, 7.5, ACCENT)
    box(slide, 0.5, 2.5, 12, 1.2, title, size=40, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)
    if subtitle:
        box(slide, 0.5, 3.9, 12, 0.8, subtitle, size=20, color=LGRAY,
            align=PP_ALIGN.CENTER)

def content_slide(slide, title, body_fn):
    bg(slide)
    rect(slide, 0, 0, 13.33, 1.05, PANEL)
    rect(slide, 0, 1.02, 13.33, 0.05, ACCENT)
    box(slide, 0.45, 0.18, 12.4, 0.75, title, size=26, bold=True, color=WHITE)
    body_fn(slide)

# ── Diagram node: rect + label (+ optional sub-label) ────────────────────────
# All text stays strictly inside; always pass explicit heights.
def node(slide, l, t, w, h, label, sub="", fill=ACCENT, tc=BG, lsize=13, ssize=10):
    rect(slide, l, t, w, h, fill)
    if sub:
        lh = h * 0.50          # top half for label
        sh = h * 0.42          # bottom half for sub (leaves a little margin)
        box(slide, l+0.06, t+0.05,    w-0.12, lh,   label, size=lsize, bold=True,
            color=tc, align=PP_ALIGN.CENTER)
        box(slide, l+0.06, t+lh+0.06, w-0.12, sh,   sub,   size=ssize,
            color=tc, align=PP_ALIGN.CENTER, italic=True)
    else:
        box(slide, l+0.06, t+0.06, w-0.12, h-0.12, label, size=lsize, bold=True,
            color=tc, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK); bg(sl)
rect(sl, 0, 0, 0.10, 7.5, ACCENT)
box(sl, 0.6, 0.8,  12.1, 0.7,
    "Department of Modeling & Simulation  ·  Old Dominion University",
    size=14, color=LGRAY)
box(sl, 0.6, 1.55, 12.1, 1.7,
    "Learning LLVM Pass Sequences\nvia Reinforcement Learning",
    size=38, bold=True, color=WHITE)
box(sl, 0.6, 3.2,  12.1, 0.55, "with Autoregressive Policies",
    size=24, color=ACCENT2)
box(sl, 0.6, 4.0,  6,    0.45, "Evan Black", size=18, bold=True, color=WHITE)
box(sl, 0.6, 4.5,  6,    0.4,  "eblac013@odu.edu", size=15, color=LGRAY)
box(sl, 0.6, 6.8,  12.1, 0.4,
    "Reinforcement Learning  ·  LLVM  ·  Pass Ordering  ·  PPO  ·  Transformer",
    size=13, color=ACCENT, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# 2 — LLVM & IR Background  (NEW)
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
def _s2_llvm(sl):
    # ── Compilation pipeline (top band) ──────────────────────────────────────
    stages = [
        ("C Source",       ".c file\nplatform-agnostic",          PANEL,  ACCENT),
        ("LLVM IR",        "typed SSA\nplatform-independent",     PANEL,  YELLOW),
        ("Optimized IR",   "pass pipeline\napplied by opt",       PANEL,  ORANGE),
        ("Native Binary",  "machine code\nfor target arch",       PANEL,  ACCENT2),
    ]
    bw = 2.50; bh = 1.15; gap = 0.30; sx = 0.35; ty = 1.22
    for i, (title, sub, bg_c, hdr_c) in enumerate(stages):
        x = sx + i*(bw+gap)
        rect(sl, x, ty, bw, bh, bg_c)
        rect(sl, x, ty, bw, 0.34, hdr_c)
        box(sl, x+0.06, ty+0.04, bw-0.12, 0.27, title,
            size=13, bold=True, color=BG, align=PP_ALIGN.CENTER)
        box(sl, x+0.06, ty+0.40, bw-0.12, 0.70, sub,
            size=11, color=LGRAY, align=PP_ALIGN.CENTER, italic=True)
        if i < len(stages)-1:
            box(sl, x+bw+0.03, ty+0.42, gap-0.06, 0.34, "→",
                size=20, color=LGRAY, align=PP_ALIGN.CENTER)
    # Commands below each stage
    cmds = [
        "clang -O3\n-Xclang -disable-llvm-optzns\n-emit-llvm -S",
        "opt --passes=\n\"inline,sroa,mem2reg,...\"",
        "clang -O3\n-Xclang -disable-llvm-passes",
        "",
    ]
    for i, cmd in enumerate(cmds):
        if not cmd: continue
        x = sx + i*(bw+gap)
        box(sl, x+0.05, ty+bh+0.06, bw-0.1, 0.62, cmd,
            size=9, color=RGBColor(0x88,0xAA,0xCC), align=PP_ALIGN.CENTER, italic=True)

    # ── What is IR? (left bottom) ─────────────────────────────────────────────
    rect(sl, 0.35, 3.06, 5.55, 3.82, PANEL)
    box(sl, 0.45, 3.12, 5.35, 0.36, "What is LLVM IR?",
        size=16, bold=True, color=ACCENT)
    bullet_box(sl, 0.45, 3.54, 5.35, 3.2, [
        (0, "Platform-independent intermediate representation"),
        (0, "Static Single Assignment (SSA) form:"),
        (1, "Every variable defined exactly once"),
        (1, "phi nodes merge values from different paths"),
        (0, "Typed: i32, float, ptr, vectors, structs"),
        (0, "Explicit control-flow graph (basic blocks + branches)"),
        (0, "Close to machine code but still target-agnostic"),
        (0, "Right level for optimization:"),
        (1, "High enough to reason about semantics"),
        (1, "Low enough to see memory, loops, calls directly"),
    ], size=13, color=LGRAY, rich=False)

    # ── IR snippet (middle bottom) ────────────────────────────────────────────
    rect(sl, 6.1, 3.06, 4.3, 3.82, RGBColor(0x0A, 0x10, 0x1C))
    box(sl, 6.2, 3.12, 4.1, 0.32, "Example: array_sum  (before passes)",
        size=12, bold=True, color=ACCENT2)
    snippet = (
        "define i32 @sum(ptr %arr, i32 %n) {\n"
        "entry:\n"
        "  %res = alloca i32          ; stack slot\n"
        "  store i32 0, ptr %res\n"
        "  br label %loop\n"
        "loop:\n"
        "  %i = phi i32 [0,%entry],[%i2,%loop]\n"
        "  %p = getelementptr i32,ptr %arr, i32 %i\n"
        "  %v = load i32, ptr %p\n"
        "  %s = load i32, ptr %res\n"
        "  %s2 = add i32 %s, %v\n"
        "  store i32 %s2, ptr %res    ; mem2reg promotes\n"
        "  %i2 = add i32 %i, 1\n"
        "  %done = icmp eq i32 %i2, %n\n"
        "  br i1 %done, label %exit, label %loop\n"
        "exit:\n"
        "  %r = load i32, ptr %res\n"
        "  ret i32 %r\n"
        "}"
    )
    box(sl, 6.15, 3.50, 4.2, 3.28, snippet,
        size=8, color=RGBColor(0xA8, 0xD8, 0xA8),
        align=PP_ALIGN.LEFT, italic=False)

    # ── Why passes matter (right bottom) ─────────────────────────────────────
    rect(sl, 10.6, 3.06, 2.55, 3.82, PANEL)
    box(sl, 10.7, 3.12, 2.35, 0.32, "After mem2reg + instcombine:",
        size=12, bold=True, color=YELLOW)
    after = (
        "define i32 @sum(ptr %arr,\n"
        "               i32 %n) {\n"
        "entry:\n"
        "  br label %loop\n"
        "loop:\n"
        "  %i = phi i32 [0,%e],[%i2,%l]\n"
        "  %s = phi i32 [0,%e],[%s2,%l]\n"
        "  %p = gep i32,ptr %arr,%i\n"
        "  %v = load i32, ptr %p\n"
        "  %s2 = add i32 %s, %v\n"
        "  %i2 = add i32 %i, 1\n"
        "  %ok = icmp ne i32 %i2, %n\n"
        "  br i1 %ok, %l, %exit\n"
        "exit:\n"
        "  ret i32 %s2\n"
        "}"
    )
    box(sl, 10.65, 3.50, 2.42, 3.28, after,
        size=8, color=RGBColor(0xA8, 0xD8, 0xFF),
        align=PP_ALIGN.LEFT, italic=False)
content_slide(sl, "Background: LLVM, IR, and Optimization Passes", _s2_llvm)


# ═══════════════════════════════════════════════════════════════════════════════
# 3 — Motivation
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
def _s3(sl):
    bullet_box(sl, 0.45, 1.25, 7.0, 5.6, [
        (0, "LLVM exposes hundreds of individual optimization passes"),
        (1, "inlining, vectorization, loop unrolling, dead-code elim, …"),
        (0, "Order matters: inline before const-prop → more constants to fold"),
        (0, "-O3 uses a fixed pipeline tuned for the average program"),
        (1, "Suboptimal for any specific program"),
        (0, "Finding the best pipeline is NP-hard (phase-ordering problem)"),
        (0, "ML offers a principled alternative:"),
        (1, "Map IR features → pass choices"),
        (1, "Train on real execution-time feedback"),
        (1, "Generalise to unseen programs"),
    ], size=17, rich=True)
    rect(sl, 7.7, 1.3, 5.2, 5.5, PANEL)
    box(sl, 7.9, 1.5, 4.9, 0.45, "The Goal", size=17, bold=True, color=ACCENT2)
    box(sl, 7.9, 2.0, 4.9, 1.2,
        "Beat -O3 on individual programs by learning a per-program pass sequence.",
        size=16, color=WHITE)
    box(sl, 7.9, 3.4, 4.9, 0.45, "Prior Work", size=17, bold=True, color=ACCENT2)
    bullet_box(sl, 7.9, 3.9, 4.9, 2.6, [
        (0, "AutoPhase — LSTM + PPO, targets IR size"),
        (0, "CompilerGym — standard RL env for LLVM"),
        (0, "MLGO — learned inlining at Google scale"),
        (0, "MILEPOST — nearest-neighbour on static features"),
        (1, "Our focus: wall-clock speedup, offline corpus + autoregressive policy"),
    ], size=14, color=LGRAY)
content_slide(sl, "Motivation & Problem", _s3)


# ═══════════════════════════════════════════════════════════════════════════════
# 4 — MDP Formulation  (State card: 2-line opcodes so metadata stays on screen)
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
def _s4(sl):
    # State card
    rect(sl, 0.35, 1.2, 3.95, 5.6, PANEL)
    rect(sl, 0.35, 1.2, 3.95, 0.50, ACCENT)
    box(sl, 0.45, 1.25, 3.75, 0.40, "State  sₜ", size=19, bold=True, color=BG)
    box(sl, 0.45, 1.80, 3.75, 0.36, "48-dim IR delta vector  fₜ ∈ ℝ⁴⁸",
        size=13, bold=True, color=YELLOW)
    box(sl, 0.45, 2.20, 3.75, 0.30, "3 chunk-pairs × 16 dims each:",
        size=12, color=LGRAY, italic=True)
    box(sl, 0.45, 2.55, 3.75, 0.28, "Opcode categories (12):",
        size=12, bold=True, color=ACCENT2)
    box(sl, 0.50, 2.87, 3.65, 0.34,
        "memory · int-arith · float-arith · bitwise · compare · cast",
        size=10, color=LGRAY, italic=True)
    box(sl, 0.50, 3.24, 3.65, 0.34,
        "call · control · phi/select · vector · block-boundary · other",
        size=10, color=LGRAY, italic=True)
    box(sl, 0.45, 3.65, 3.75, 0.28, "Metadata rates (4):",
        size=12, bold=True, color=ACCENT2)
    box(sl, 0.50, 3.97, 3.65, 0.28,
        "tbaa · llvm.loop · alias.scope · noalias",
        size=10, color=LGRAY, italic=True)
    rect(sl, 0.45, 4.32, 3.75, 0.03, RGBColor(0x33, 0x44, 0x66))
    box(sl, 0.45, 4.40, 3.75, 0.36,
        "Δ = chunk(c+1) − chunk(c)  for c ∈ {0,1,2}",
        size=11, color=WHITE)
    box(sl, 0.45, 4.82, 3.75, 0.52,
        "Invariant to function size; responds to\nloop-rotate, licm, sroa structural shifts",
        size=11, color=LGRAY, italic=True)

    # Action card
    rect(sl, 4.50, 1.2, 4.20, 5.6, PANEL)
    rect(sl, 4.50, 1.2, 4.20, 0.50, ACCENT2)
    box(sl, 4.60, 1.25, 4.00, 0.40, "Action  aₜ", size=19, bold=True, color=BG)
    box(sl, 4.60, 1.80, 4.00, 0.36, "29 total:  28 passes  +  1 Stop",
        size=13, bold=True, color=YELLOW)
    box(sl, 4.60, 2.22, 4.00, 0.28, "Example passes:",
        size=12, bold=True, color=ACCENT)
    for i, (nm, desc) in enumerate([
        ("inline",       "removes call overhead; exposes callee"),
        ("sroa",         "scalar replace aggregates → registers"),
        ("mem2reg",      "alloca → SSA  (prerequisite for most)"),
        ("instcombine",  "fold/simplify instruction patterns"),
        ("simplifycfg",  "clean dead blocks, merge branches"),
        ("loop-rotate",  "canonicalize loops for LICM"),
        ("licm",         "hoist loop-invariant code out"),
        ("loop-unroll",  "replicate loop body N times"),
        ("gvn",          "global value numbering / CSE"),
        ("Stop",         "terminate early → compile binary"),
    ]):
        y = 2.56 + i*0.29
        box(sl, 4.60,  y, 1.22, 0.28, nm,   size=11, bold=True, color=ACCENT2)
        box(sl, 5.85,  y, 2.75, 0.28, desc, size=11, color=LGRAY, italic=True)

    # Reward card
    rect(sl, 8.90, 1.2, 4.08, 5.6, PANEL)
    rect(sl, 8.90, 1.2, 4.08, 0.50, YELLOW)
    box(sl, 9.00, 1.25, 3.88, 0.40, "Reward  r", size=19, bold=True, color=BG)
    box(sl, 9.00, 1.80, 3.88, 0.36, "Terminal speedup vs -O3",
        size=13, bold=True, color=WHITE)
    box(sl, 9.00, 2.22, 3.88, 0.40, "r = (t_O3 − t_policy) / t_O3",
        size=14, bold=True, color=YELLOW)
    box(sl, 9.00, 2.68, 3.88, 0.28, "Positive = faster than -O3",
        size=12, color=ACCENT2, italic=True)
    rect(sl, 9.00, 3.04, 3.78, 3.08, RGBColor(0x0A, 0x10, 0x1C))
    box(sl, 9.10, 3.10, 3.60, 0.32, "Concrete example — fft:",
        size=12, bold=True, color=ACCENT)
    box(sl, 9.10, 3.48, 3.60, 0.30, "t_O3  = 46,509 ns",  size=13, color=LGRAY)
    box(sl, 9.10, 3.82, 3.60, 0.30, "t_policy = 41,852 ns", size=13, color=LGRAY)
    rect(sl, 9.10, 4.18, 3.58, 0.03, LGRAY)
    box(sl, 9.10, 4.26, 3.60, 0.34,
        "r = (46,509 − 41,852) / 46,509", size=12, color=WHITE)
    box(sl, 9.10, 4.64, 3.60, 0.34, "  = 4,657 / 46,509",  size=12, color=WHITE)
    box(sl, 9.10, 5.02, 3.60, 0.42, "  = +10.0%  speedup",
        size=16, bold=True, color=ACCENT2)
    box(sl, 9.00, 5.52, 3.88, 0.40,
        "Trimmed mean · 201 iterations · 5 warm-up\nCLOCK_MONOTONIC",
        size=10, color=LGRAY, italic=True)
content_slide(sl, "MDP Formulation", _s4)


# ═══════════════════════════════════════════════════════════════════════════════
# 5 — Section: Dataset & EDA
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
section_header(sl, "Dataset & Exploratory Data Analysis",
               "950k random sequences · 38 benchmarks · bench-cache")


# ═══════════════════════════════════════════════════════════════════════════════
# 6 — Dataset Landscape
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
def _s6(sl):
    rect(sl, 0.35, 1.2, 4.3, 5.95, PANEL)
    box(sl, 0.45, 1.27, 4.1, 0.38, "Bench-Cache Statistics",
        size=15, bold=True, color=ACCENT)
    stats = [
        ("Total sequences",    "950,000"),
        ("Benchmarks in pool", "38"),
        ("Training benchmarks","6"),
        ("Sequences / fn",     "25,000"),
        ("Sequence lengths",   "1 – 20"),
        ("Pass menu size",     "28 + Stop"),
        ("Baselines per fn",   "-O0, -O2, -O3"),
        ("Cache format",       "on-disk key-value"),
    ]
    for i, (k, v) in enumerate(stats):
        y = 1.74 + i*0.60
        rect(sl, 0.45, y, 4.1, 0.56, BG if i%2==0 else RGBColor(0x12,0x1D,0x35))
        box(sl, 0.55, y+0.09, 2.1, 0.38, k,  size=12, color=LGRAY)
        box(sl, 2.65, y+0.09, 1.8, 0.38, v,  size=12, bold=True, color=YELLOW)
    cg = os.path.join(CKPT, "ceiling_gaps.png")
    sl_img = os.path.join(CKPT, "seq_length_by_tier.png")
    if os.path.exists(cg):
        img(sl, cg,    4.75, 1.2,  4.3)
    if os.path.exists(sl_img):
        img(sl, sl_img, 9.15, 1.2, 4.0)
content_slide(sl, "Training Landscape — Bench-Cache Overview", _s6)


# ═══════════════════════════════════════════════════════════════════════════════
# 7 — Benchmark Pool
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
def _s7(sl):
    bullet_box(sl, 0.45, 1.25, 5.55, 2.05, [
        (0, "Pool: 38 self-contained C benchmarks"),
        (1, "Sorting, numerical, data structures, codecs, recursive"),
        (0, "25,000 random sequences per benchmark → ~950k entries"),
        (0, "3 baselines (-O0, -O2, -O3) recorded per benchmark"),
    ], size=16, rich=True)
    box(sl, 0.45, 3.38, 5.55, 0.40, "6 Training Benchmarks — selection criteria:",
        size=15, bold=True, color=ACCENT2)
    bullet_box(sl, 0.45, 3.84, 5.55, 2.5, [
        (0, "At least one random sequence beats -O3"),
        (0, "High CV across sequences (rugged reward landscape)"),
        (0, "Diverse computational profiles"),
    ], size=15, rich=True)
    spd_vals = [27.2, 24.3, 12.7, 10.0, 4.7, 9.2]
    cv_vals  = [13.9, 42.7, 46.0, 35.5, 36.6, 18.4]
    cold_g=(0x66,0xDD,0x88); hot_g=(0x00,0xFF,0x66)
    cold_o=(0xCC,0xAA,0x44); hot_o=(0xFF,0xD7,0x00)
    rows = [
        ("Benchmark",      "O3 (ns)",  "Best (ns)","vs O3",  "CV"),
        ("interpreter",    "37,918",   "27,591",   "−27.2%", "13.9%"),
        ("kmp_search",     "886,408",  "671,195",  "−24.3%", "42.7%"),
        ("polynomial_eval","417,235",  "364,235",  "−12.7%", "46.0%"),
        ("fft",            "46,509",   "41,852",   "−10.0%", "35.5%"),
        ("array_reduction","97,184",   "92,619",   " −4.7%", "36.6%"),
        ("binary_tree",    "117,990",  "107,087",  " −9.2%", "18.4%"),
    ]
    col_x=[6.3,7.9,9.3,10.7,11.9]; col_w=[1.55,1.35,1.35,1.15,1.1]
    for ri, row in enumerate(rows):
        y = 1.25 + ri*0.83
        if ri == 0:
            rect(sl, 6.25, y, 6.75, 0.8, ACCENT)
            for cell,cx,cw in zip(row,col_x,col_w):
                box(sl,cx,y+0.12,cw,0.56,cell,size=13,bold=True,color=BG,align=PP_ALIGN.CENTER)
        else:
            di=ri-1
            rect(sl,6.25,y,6.75,0.8,PANEL if ri%2==1 else BG)
            sc=lerp_rgb(cold_g,hot_g,(spd_vals[di]-4.7)/(27.2-4.7))
            cc=lerp_rgb(cold_o,hot_o,(cv_vals[di]-13.9)/(46.0-13.9))
            for ci,(cell,cx,cw) in enumerate(zip(row,col_x,col_w)):
                fc=WHITE if ci==0 else (sc if ci==3 else (cc if ci==4 else LGRAY))
                box(sl,cx,y+0.12,cw,0.56,cell,size=13,bold=(ci in(3,4)),
                    color=fc,align=PP_ALIGN.CENTER)
    box(sl,6.25,6.55,3.3,0.40,"Green = speedup magnitude",size=12,color=ACCENT2,italic=True)
    box(sl,9.60,6.55,3.4,0.40,"Yellow = landscape ruggedness (CV)",size=12,color=YELLOW,italic=True)
content_slide(sl, "Benchmark Pool & Training Benchmark Selection", _s7)


# ═══════════════════════════════════════════════════════════════════════════════
# 8 — IR Features
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
def _s8(sl):
    bullet_box(sl, 0.45, 1.25, 5.8, 3.5, [
        (0, "Policy observes a 48-dim IR feature vector fₜ"),
        (0, "Split IR into C=4 equal positional chunks"),
        (0, "Per chunk: opcode histogram (12 categories) + metadata rates (4 kinds)"),
        (1, "categories: memory, int-arith, float-arith, bitwise, compare, cast,"),
        (1, "call, control, phi/select, vector, block-boundary, other"),
        (1, "metadata: tbaa, llvm.loop, alias.scope, noalias"),
        (0, "3 adjacent-chunk deltas → 3 × 16 = 48-dim vector"),
        (0, "Key properties:"),
        (1, "Invariant to function size"),
        (1, "Responds to loop-rotate, licm, sroa — invisible to global counts"),
    ], size=15, rich=True)
    img(sl, os.path.join(BASE,"features_op.png"),   6.3, 1.25, 6.7)
    img(sl, os.path.join(BASE,"features_meta.png"), 6.3, 4.3,  6.7)
content_slide(sl, "IR Feature Representation", _s8)


# ═══════════════════════════════════════════════════════════════════════════════
# 9 — Pass Menu & EDA
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
def _s9(sl):
    bullet_box(sl, 0.45, 1.25, 5.9, 2.2, [
        (0, "28-pass primary menu + Stop token (29 total)"),
        (0, "Inclusion criteria:"),
        (1, "Appears in -O3 inner loop or is a prerequisite"),
        (1, "Demonstrable impact on C compute code"),
        (1, "Covers a distinct optimization category"),
        (0, "LLVM 20 nesting handled transparently:"),
        (1, "inline → cgscc(inline)"),
        (1, "licm → loop-mssa(licm)  ·  loop-rotate → loop(loop-rotate)"),
    ], size=15, rich=True)
    box(sl, 0.45, 3.55, 5.9, 0.38, "Top passes by EDA analysis:",
        size=15, bold=True, color=ACCENT2)
    rows2=[
        ("Pass",        "Top-Decile Enrich.","Geo-Mean Speedup"),
        ("inline",      "2.35×",             "1.26×"),
        ("simplifycfg", "1.98×",             "1.12×"),
        ("sroa",        "1.87×",             "1.09×"),
        ("mem2reg",     "1.76×",             "1.08×"),
    ]
    for ri,row in enumerate(rows2):
        y=4.00+ri*0.52
        rect(sl,0.45,y,5.9,0.49,ACCENT if ri==0 else (PANEL if ri%2 else BG))
        for cell,cx,cw in zip(row,[0.5,2.05,3.9],[1.5,1.8,2.0]):
            fc=BG if ri==0 else (YELLOW if cx>2 else WHITE)
            box(sl,cx,y+0.05,cw,0.39,cell,size=13,bold=(ri==0),color=fc)
    bullet_box(sl, 6.5, 1.25, 6.5, 5.7, [
        (0, "inline is the top pass-enabler:"),
        (1, "Enrichment 2.35× = in top-decile sequences at 2.35× baseline rate"),
        (1, "Exposes callee bodies to const-prop, DCE, loop transforms"),
        (0, "simplifycfg + sroa + mem2reg unlock downstream analyses"),
        (1, "SSA promotion (mem2reg) required before most scalar optimizations"),
        (0, "Sequence length distribution ~uniform across 1–20"),
        (0, "Best-found median length: 17–18 (no cap spike → K=20 sufficient)"),
        (0, "Bench-cache lookup rate grows during training:"),
        (1, "Policy re-generates previously benchmarked sequences → instant reward"),
    ], size=15, rich=True)
content_slide(sl, "Pass Menu & EDA Highlights", _s9)


# ═══════════════════════════════════════════════════════════════════════════════
# 10 — Section: Architecture & Training
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
section_header(sl, "Policy Architecture & Training",
               "Auto-TFX  ·  Auto-GRU  ·  PPO")


# ═══════════════════════════════════════════════════════════════════════════════
# 11 — Architecture Diagrams  (rebuilt with pixel-precise positions)
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
def _s11(sl):
    # ── Auto-TFX  (x: 0.30 → 6.38) ──────────────────────────────────────────
    rect(sl, 0.30, 1.20, 6.10, 0.46, ACCENT)
    box(sl, 0.34, 1.24, 6.02, 0.36,
        "Auto-TFX  (Causal Transformer)", size=15, bold=True,
        color=BG, align=PP_ALIGN.CENTER)

    # Token row  y=1.76 h=0.90
    tok = [
        ("W·fₜ",     "IR proj\n48→128",  ACCENT),
        ("E[a₀]",    "action\nembed",     RGBColor(0x3A,0x8E,0xEF)),
        ("E[a₁]",    "action\nembed",     RGBColor(0x2A,0x7E,0xDF)),
        ("…",        "",                  PANEL),
        ("E[aₜ₋₁]",  "action\nembed",    RGBColor(0x1A,0x6E,0xCF)),
    ]
    TW=1.07; TH=0.90
    for i,(lbl,sub,fc) in enumerate(tok):
        node(sl, 0.32+i*1.14, 1.76, TW, TH, lbl, sub,
             fill=fc, tc=WHITE if fc!=PANEL else LGRAY, lsize=12, ssize=9)

    box(sl, 1.60, 2.72, 3.0, 0.28,
        "↓  full cross-step attention  ↓", size=11, color=LGRAY, align=PP_ALIGN.CENTER)

    # Transformer  y=3.04 h=0.80
    node(sl, 0.32, 3.04, 5.88, 0.80,
         "Transformer Encoder",
         "2 layers · 4 heads · d=128 · FF 256 · dropout 0.1",
         fill=RGBColor(0x1E,0x40,0x7A), tc=WHITE, lsize=14, ssize=11)

    box(sl, 0.72, 3.90, 2.0, 0.26, "↓  token 0", size=11, color=LGRAY, align=PP_ALIGN.CENTER)
    box(sl, 3.68, 3.90, 2.0, 0.26, "↓  token t", size=11, color=LGRAY, align=PP_ALIGN.CENTER)

    # Value/Policy heads  y=4.20 h=0.95  — wide enough, tall enough
    VH=0.95
    node(sl, 0.32, 4.20, 2.68, VH, "Value Head", "reads token 0\n→  V(sₜ)",
         fill=ACCENT2, tc=BG, lsize=13, ssize=11)
    node(sl, 3.38, 4.20, 2.68, VH, "Policy Head", "reads token t\n→  logits (29)",
         fill=ORANGE,  tc=BG, lsize=13, ssize=11)

    rect(sl, 0.32, 5.22, 5.88, 0.74, RGBColor(0x0A,0x10,0x1C))
    bullet_box(sl, 0.40, 5.25, 5.72, 0.68, [
        (0, "No causal mask — decode 1 step at a time; full sequence always visible"),
        (0, "Batched replay: K forward passes (one per step position) vs N×K naïve"),
    ], size=11, color=LGRAY)

    # ── Auto-GRU  (x: 6.92 → 13.0) ──────────────────────────────────────────
    rect(sl, 6.92, 1.20, 6.10, 0.46, ACCENT2)
    box(sl, 6.96, 1.24, 6.02, 0.36,
        "Auto-GRU  (Recurrent)", size=15, bold=True,
        color=BG, align=PP_ALIGN.CENTER)

    # Init row  y=2.02 h=0.84
    IH=0.84
    box(sl, 6.96, 1.74, 5.9, 0.24, "Step 0 — seed hidden state from initial IR:",
        size=11, bold=True, color=ACCENT)
    node(sl, 6.96, 2.02, 1.22, IH,  "f₀",  "IR feat",  fill=ACCENT, tc=BG, lsize=14, ssize=10)
    box(sl, 8.24, 2.38, 0.30, 0.24, "→", size=14, color=LGRAY, align=PP_ALIGN.CENTER)
    node(sl, 8.60, 2.02, 1.55, IH, "Linear\nproject", "", fill=RGBColor(0x1E,0x40,0x7A), tc=WHITE, lsize=12)
    box(sl, 10.21,2.38, 0.30, 0.24, "→", size=14, color=LGRAY, align=PP_ALIGN.CENTER)
    node(sl, 10.57,2.02, 2.18, IH, "h₀", "initial hidden state", fill=ACCENT2, tc=BG, lsize=15, ssize=11)

    # Recurrent section  — fₜ and aₜ₋₁ stacked on left
    box(sl, 6.96, 3.00, 5.9, 0.24, "Step t — recurrent update:",
        size=11, bold=True, color=ACCENT)
    FH=0.72  # height of input nodes
    node(sl, 6.96, 3.28, 1.22, FH, "fₜ",    "IR feat", fill=ACCENT,                      tc=BG, lsize=13, ssize=10)
    node(sl, 6.96, 4.06, 1.22, FH, "aₜ₋₁",  "action",  fill=RGBColor(0x3A,0x8E,0xEF),    tc=BG, lsize=13, ssize=10)
    # fₜ ends y=4.00, aₜ₋₁ y=4.06-4.78

    # GRU cell spans both inputs  y=3.28 h=1.50
    node(sl, 8.55, 3.28, 1.55, 1.50, "GRU\nCell", "",
         fill=RGBColor(0x1E,0x40,0x7A), tc=WHITE, lsize=13)
    box(sl, 8.22, 3.55, 0.30, 0.24, "→", size=13, color=LGRAY, align=PP_ALIGN.CENTER)
    box(sl, 8.22, 4.22, 0.30, 0.24, "→", size=13, color=LGRAY, align=PP_ALIGN.CENTER)
    box(sl, 10.16,4.00, 0.30, 0.24, "→", size=13, color=LGRAY, align=PP_ALIGN.CENTER)

    # hₜ  centered at y=3.78 (GRU center = 3.28+0.75=4.03 → hₜ center=4.03)
    node(sl, 10.52, 3.66, 2.14, 0.80, "hₜ", "hidden state", fill=ACCENT2, tc=BG, lsize=15, ssize=11)
    # recurrent label below GRU
    box(sl, 8.55, 4.84, 1.55, 0.26, "↑ hₜ₋₁  (recurrent)",
        size=10, color=ORANGE, align=PP_ALIGN.CENTER, italic=True)

    # Output heads  y=5.10 h=0.92  — starts 0.32 below aₜ₋₁ end (4.78)
    box(sl, 6.96, 5.08, 5.9, 0.24, "→  output heads:",
        size=11, color=LGRAY)
    OH=0.92
    node(sl, 6.96, 5.36, 2.68, OH, "Value Head", "→  V(sₜ)",
         fill=ACCENT2, tc=BG, lsize=13, ssize=12)
    node(sl, 9.80, 5.36, 2.68, OH, "Policy Head", "→  logits (29)",
         fill=ORANGE,  tc=BG, lsize=13, ssize=12)

    rect(sl, 6.96, 6.38, 5.88, 0.48, RGBColor(0x0A,0x10,0x1C))
    bullet_box(sl, 7.04, 6.40, 5.72, 0.44, [
        (0, "Recency bias: recent IR state has more influence — free inductive prior"),
        (0, "Compresses full episode history into a fixed-width hidden vector"),
    ], size=11, color=LGRAY)

    box(sl, 0.30, 6.94, 12.73, 0.36,
        "Both architectures re-encode current IR features fₜ at every step — neither operates on a stale snapshot.",
        size=12, color=LGRAY, italic=True, align=PP_ALIGN.CENTER)
content_slide(sl, "Policy Architecture Diagrams", _s11)


# ═══════════════════════════════════════════════════════════════════════════════
# 12 — Pass Selection Walkthrough
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
def _s12_walk(sl):
    rect(sl, 0.35, 1.20, 12.6, 0.68, PANEL)
    rect(sl, 0.35, 1.20, 12.6, 0.07, YELLOW)
    box(sl, 0.45, 1.32, 12.4, 0.44,
        "Example: interpreter  ·  -O3 baseline = 37,275 ns  ·  "
        "Greedy policy result: ~27,183 ns  →  +27.1% speedup",
        size=14, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
    steps = [
        ("Step 0",  "fₜ: mem↑\ncall↑\nctrl↑",   "inline",            "remove call\noverhead",          ACCENT),
        ("Step 1",  "fₜ: mem↑\ncfg↑\nint↓",     "simplifycfg",       "clean dead\nblocks/branches",    RGBColor(0x3A,0x8E,0xEF)),
        ("Step 2",  "fₜ: mem↓↓\nint↑\nSSA↑",    "sroa",              "stack→regs\n(mem↓ heavy)",        RGBColor(0x2A,0xAE,0x7F)),
        ("Step 3",  "fₜ: mem≈0\nconsts\nvisible","mem2reg",           "alloca→SSA\npromotes vals",       RGBColor(0x5A,0x9E,0x4F)),
        ("Step 4",  "fₜ: stable\nΔ small\n/step","instcombine",       "fold/simplify\npatterns",         RGBColor(0x9A,0x7E,0x2F)),
        ("5–19",    "fₜ: minor\nΔ /step\n(safe)","instcombine\n×15", "safe: rarely\nhurts, occ. wins", RGBColor(0x55,0x55,0x77)),
    ]
    bw=1.95; gap=0.12; sx=0.35; ty=2.06
    for i,(step,feat,action,effect,col) in enumerate(steps):
        x=sx+i*(bw+gap)
        rect(sl,x,ty,bw,0.38,col)
        box(sl,x+0.05,ty+0.05,bw-0.1,0.28,step,size=12,bold=True,color=BG,align=PP_ALIGN.CENTER)
        rect(sl,x,ty+0.41,bw,0.76,RGBColor(0x0A,0x10,0x1C))
        box(sl,x+0.05,ty+0.46,bw-0.1,0.66,feat,size=10,color=LGRAY,align=PP_ALIGN.CENTER,italic=True)
        box(sl,x+0.30,ty+1.22,bw-0.60,0.24,"↓ picks",size=10,color=LGRAY,align=PP_ALIGN.CENTER)
        rect(sl,x,ty+1.50,bw,0.52,col)
        box(sl,x+0.04,ty+1.55,bw-0.08,0.42,action,size=11,bold=True,color=BG,align=PP_ALIGN.CENTER)
        rect(sl,x,ty+2.06,bw,0.68,PANEL)
        box(sl,x+0.04,ty+2.12,bw-0.08,0.58,effect,size=10,color=WHITE,align=PP_ALIGN.CENTER,italic=True)
        if i<len(steps)-1:
            box(sl,x+bw+0.01,ty+1.66,gap+0.09,0.28,"→",size=14,color=LGRAY,align=PP_ALIGN.CENTER)
    rect(sl,0.35,ty+2.87,12.6,0.54,RGBColor(0x0A,0x22,0x14))
    rect(sl,0.35,ty+2.87,12.6,0.06,ACCENT2)
    box(sl,0.45,ty+2.96,12.4,0.38,
        "Compile → benchmark  ·  measured = 27,183 ns  ·  "
        "r = (37,275 − 27,183) / 37,275 = +27.1%  ✓",
        size=14,bold=True,color=ACCENT2,align=PP_ALIGN.CENTER)
    rect(sl,0.35,ty+3.54,12.6,1.06,PANEL)
    box(sl,0.45,ty+3.60,5.9,0.32,"Why repeat instcombine?",size=13,bold=True,color=YELLOW)
    bullet_box(sl,0.45,ty+3.96,5.85,0.60,[
        (0,"Q(Stop)=current speedup; Q(pass)=E[future] — uncertain, possibly higher"),
        (0,"Penalty p=0.025 per no-op is tiny vs ~28% gain — rational to continue"),
    ],size=11,color=LGRAY,rich=True)
    box(sl,6.60,ty+3.60,6.2,0.32,"Why is inline so dominant?",size=13,bold=True,color=ACCENT2)
    bullet_box(sl,6.60,ty+3.96,6.15,0.60,[
        (0,"interpreter has a tight bytecode dispatch loop — call overhead dominates"),
        (0,"After inline: callee body visible → const-prop + DCE fire hard"),
    ],size=11,color=LGRAY)
content_slide(sl, "Concrete Pass-Selection Walkthrough", _s12_walk)


# ═══════════════════════════════════════════════════════════════════════════════
# 13 — PPO Training  (formula expanded into 3 colour-coded panels)
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
def _s13(sl):
    # Left: overview + hyperparams
    bullet_box(sl, 0.45, 1.25, 5.65, 2.90, [
        (0, "On-policy collect → update loop"),
        (0, "Collect 256 episodes/epoch (Rayon parallel)"),
        (0, "4 PPO inner epochs on mini-batches of 128"),
        (0, "KL early-stop ≤ 0.05 aborts inner loop if update too large"),
        (0, "AdamW · cosine LR 1e-3 → ~0 over 100 epochs · 100 epochs total"),
    ], size=15, rich=True)
    hparams=[
        ("d_model","128"), ("Layers","2"),   ("Heads","4"),
        ("FF dim","256"),  ("Dropout","0.1"),("Horizon K","20"),
        ("Actions","29"),  ("Ep/epoch","256"),("Mini-batch","128"),
    ]
    box(sl, 0.45, 4.28, 5.65, 0.36, "Hyperparameters", size=14, bold=True, color=ACCENT)
    for i,(k,v) in enumerate(hparams):
        col=i%3; row=i//3
        x=0.45+col*1.88; y=4.72+row*0.50
        rect(sl,x,y,1.85,0.47,PANEL if row%2==0 else BG)
        box(sl,x+0.07,y+0.05,1.0,0.37,k,size=11,color=LGRAY)
        box(sl,x+1.07,y+0.05,0.75,0.37,v,size=11,bold=True,color=YELLOW)

    # Right: full objective label + 3 panels
    box(sl, 6.35, 1.22, 6.65, 0.34,
        "L  =  −E[ L_clip ]  +  c_v · E[ L_value ]  −  c_e · E[ L_entropy ]",
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    PH = 1.78   # panel height — enough for title + formula + 3 rows
    panels = [
        (ACCENT,
         "CLIP Term  — policy gradient",
         "L_clip = min( ρₜ · Âₜ ,  clip(ρₜ, 1−ε, 1+ε) · Âₜ )",
         [("ρₜ = π_θ(aₜ|sₜ) / π_old(aₜ|sₜ)",
           "probability ratio new policy over old collection policy"),
          ("Âₜ = Rₜ − V_θ(sₜ)",
           "advantage: return minus value baseline, batch-normalised"),
          ("ε = 0.2",
           "clip threshold — hard limit on policy drift per gradient step")]),
        (ORANGE,
         "Value Loss  —  c_v = 0.5",
         "L_value = ( V_θ(sₜ) − Rₜ )²",
         [("V_θ(sₜ)",
           "value head prediction (token 0 in TFX, hidden state hₜ in GRU)"),
          ("Rₜ",
           "per-step return (episode or instruction-weighted formulation)"),
          ("EV → 0.69/0.59",
           "TFX/GRU episode EV at epoch 100; weighted runs reach 0.87/0.86")]),
        (ACCENT2,
         "Entropy Bonus  —  c_e = 0.03",
         "L_entropy = H( π_θ(·|sₜ) )   [Shannon entropy]",
         [("H(π)",
           "spread over all 29 actions — prevents premature convergence"),
          ("c_e = 0.03",
           "small enough not to dominate; large enough to prevent collapse"),
          ("KL ≤ 0.05",
           "per-mini-batch early-stop guard as additional safety net")]),
    ]
    for pi,(col,title,formula,rows) in enumerate(panels):
        y = 1.62 + pi*(PH+0.05)
        rect(sl, 6.35, y, 6.65, PH, PANEL)
        rect(sl, 6.35, y, 0.08, PH, col)          # colour stripe
        box(sl, 6.50, y+0.06, 6.42, 0.30,
            title, size=13, bold=True, color=col)
        box(sl, 6.50, y+0.40, 6.42, 0.34,
            formula, size=12, bold=True, color=YELLOW)
        for ri,(sym,desc) in enumerate(rows):
            ry = y + 0.80 + ri*0.34
            box(sl, 6.52, ry, 1.55, 0.30, sym,  size=10, bold=True, color=col)
            box(sl, 8.10, ry, 4.82, 0.30, desc, size=10, color=LGRAY, italic=True)
content_slide(sl, "PPO Training Setup", _s13)


# ═══════════════════════════════════════════════════════════════════════════════
# 14 — Return Formulations
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
def _s14(sl):
    rect(sl,0.45,1.20,5.90,2.90,PANEL); rect(sl,0.45,1.20,5.90,0.50,ACCENT)
    box(sl,0.55,1.25,5.70,0.40,"Episode Return",size=19,bold=True,color=BG)
    rich_box(sl,0.55,1.75,5.70,0.50,"Rₜ = r  ∀ t",size=18,base=WHITE)
    bullet_box(sl,0.55,2.35,5.70,1.50,[
        (0,"Terminal speedup assigned uniformly to all steps"),
        (0,"Unbiased maximum-likelihood estimate"),
        (0,"No assumptions about which steps were responsible"),
    ],size=14,color=LGRAY,rich=True)
    rect(sl,6.90,1.20,6.00,2.90,PANEL); rect(sl,6.90,1.20,6.00,0.50,ACCENT2)
    box(sl,7.00,1.25,5.80,0.40,"Instruction-Weighted Return",size=19,bold=True,color=BG)
    rich_box(sl,7.00,1.75,5.80,0.60,
             "Rₜ = r·|ΔIₜ|/Σ|ΔIₛ| + 0.05·sign(−ΔIₜ) − 0.025·𝟙[noopₜ]",
             size=13,base=WHITE)
    bullet_box(sl,7.00,2.45,5.80,1.40,[
        (0,"Redistributes credit proportional to instruction reduction"),
        (0,"No-op penalty 0.025; direction bonus 0.05"),
        (0,"Sum of first term = r (total reward conserved)"),
    ],size=14,color=LGRAY,rich=True)
    rect(sl,0.45,4.25,12.40,2.15,RGBColor(0x1A,0x1A,0x2E))
    rect(sl,0.45,4.25,12.40,0.42,RGBColor(0x44,0x44,0x66))
    box(sl,0.55,4.30,12.20,0.32,"IR-Step Return  (Ablation)",size=15,bold=True,color=LGRAY)
    bullet_box(sl,0.55,4.75,12.20,1.40,[
        (0,"Per-step normalised instruction-count delta as reward — bypasses benchmarking entirely"),
        (0,"Dense, low-noise signal: verifies model architecture + PPO loop independently of sparse runtime reward"),
        (0,"Not evaluated as a policy; training results are an infrastructure verification only"),
    ],size=14,color=LGRAY,rich=True)
    box(sl,0.45,6.55,12.40,0.52,
        "No-op threshold: |ΔI| < 0.01 AND L1(Δfeatures) < 0.05  "
        "(dual threshold handles loop-rotate/licm restructuring without instruction-count change)",
        size=12,color=LGRAY,italic=True)
content_slide(sl, "Return Formulations  &  Credit Assignment", _s14)


# ═══════════════════════════════════════════════════════════════════════════════
# 15 — Section: Results
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
section_header(sl, "Results", "Auto-TFX + Episode Return")


# ═══════════════════════════════════════════════════════════════════════════════
# 16 — Training Curves
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
def _s16(sl):
    ti = os.path.join(CKPT,"auto-tfx-episode.png")
    if os.path.exists(ti): img(sl,ti,0.45,1.2,7.8)
    bullet_box(sl,8.4,1.25,4.7,5.7,[
        (0,"Mean speedup: −0.40 → +0.091  (EMA +0.080)"),
        (0,"Explained Variance reaches 0.69"),
        (1,"Value head learns accurate return predictions"),
        (1,"48-dim delta features sufficient for value function"),
        (0,"Entropy: max → ~30% of max"),
        (1,"Progressive concentration, no collapse"),
        (0,"No-op % stabilises ~30%"),
        (1,"Policy learns to prefer impactful passes"),
        (0,"interpreter and kmp_search lead per-function"),
        (1,"Consistent with high random-search ceiling"),
    ],size=15,rich=True)
content_slide(sl,"Auto-TFX + Episode Return — Training Curves",_s16)


# ═══════════════════════════════════════════════════════════════════════════════
# 17 — Evaluation Results
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
def _s17(sl):
    ei = os.path.join(BASE,"eval.png")
    if os.path.exists(ei): img(sl,ei,0.45,1.2,7.5)
    rows=[
        ("Function",      "rand mean","rand best","greedy", "samp best"),
        ("array_red.",    "−0.795",   "−0.039",   "+0.075", "+0.095"),
        ("binary_tree",   "−0.089",   "+0.062",   "+0.027", "+0.066"),
        ("fft",           "−0.547",   "−0.033",   "−0.016", "+0.037"),
        ("interpreter",   "−0.252",   "+0.010",   "+0.271", "+0.276"),
        ("kmp_search",    "−0.101",   "+0.209",   "+0.237", "+0.278"),
        ("poly_eval",     "−0.261",   "+0.120",   "+0.031", "+0.149"),
    ]
    col_x=[8.10,9.65,10.65,11.50,12.30]; col_w=[1.50,0.95,0.85,0.78,0.90]
    for ri,row in enumerate(rows):
        y=1.25+ri*0.74
        rect(sl,8.05,y,5.10,0.70,ACCENT if ri==0 else (PANEL if ri%2 else BG))
        for ci,(cell,cx,cw) in enumerate(zip(row,col_x,col_w)):
            if ri==0: fc=BG
            elif ci in(3,4): fc=ACCENT2 if cell.startswith("+") else RED
            elif ci==1 and any(cell.startswith(p) for p in("−0.5","−0.6","−0.7","−0.8","−0.9","−1.")): fc=RED
            else: fc=LGRAY
            box(sl,cx,y+0.07,cw,0.56,cell,size=12,
                bold=(ri==0 or(ri>0 and ci in(3,4))),
                color=fc,align=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER)
    rect(sl,8.05,6.45,5.10,0.72,RGBColor(0x0A,0x22,0x14))
    box(sl,8.10,6.52,5.00,0.38,
        "Greedy: 5/6 beat -O3 · mean +10.4% · sample-best all 6",
        size=14,bold=True,color=ACCENT2,align=PP_ALIGN.CENTER)
content_slide(sl,"Auto-TFX + Episode Return — Evaluation vs -O3",_s17)


# ═══════════════════════════════════════════════════════════════════════════════
# 18 — GRU Episode Training + Weighted Returns Overview
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
def _s18_gru_wt(sl):
    # Left panel: Auto-GRU episode training plot
    rect(sl,0.30,1.20,6.10,0.40,ACCENT2)
    box(sl,0.34,1.24,6.02,0.30,"Auto-GRU  +  Episode Return — Training",
        size=13,bold=True,color=BG,align=PP_ALIGN.CENTER)
    gi = os.path.join(CKPT,"auto-gru-episode.png")
    if os.path.exists(gi): img(sl,gi,0.30,1.64,6.10)
    bullet_box(sl,0.35,5.55,6.05,1.55,[
        (0,"Mean speedup −0.44 → +0.043  (EMA +0.031)"),
        (0,"EV reaches 0.59  (slower than TFX, same qualitative trajectory)"),
        (0,"No-op rate ~31% — consistent with episode-return behaviour"),
        (0,"Greedy eval mean: +11.4% — comparable to Auto-TFX (+10.4%)"),
    ],size=12,color=LGRAY,rich=True)

    # Right panel: Weighted returns both worse
    rect(sl,6.70,1.20,6.38,0.40,YELLOW)
    box(sl,6.74,1.24,6.30,0.30,"Instruction-Weighted Return — Both Architectures",
        size=13,bold=True,color=BG,align=PP_ALIGN.CENTER)
    wi = os.path.join(CKPT,"auto-tfx-weighted.png")
    if os.path.exists(wi): img(sl,wi,6.70,1.64,6.38)
    bullet_box(sl,6.75,5.55,6.28,1.55,[
        (0,"TFX weighted: speedup → +0.012 (vs +0.091 episode);  EV = 0.87"),
        (0,"GRU weighted: speedup → +0.023 (vs +0.043 episode);  EV = 0.86"),
        (0,"Higher EV but lower speedup — proxy misalignment (IR ≠ runtime)"),
        (0,"No-op rate ~40–41% — elevated vs episode runs"),
    ],size=12,color=LGRAY,rich=True)
content_slide(sl,"Training Results — GRU Episode & Weighted Returns",_s18_gru_wt)


# ═══════════════════════════════════════════════════════════════════════════════
# 19 — Weighted Return Evaluation Tables + IR-step / ir_corr
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
def _s19_wt_ir(sl):
    # Left: weighted eval tables (TFX + GRU side-by-side)
    def small_table(l, t, w, title, col, rows_data):
        rect(sl,l,t,w,0.35,col)
        box(sl,l+0.04,t+0.04,w-0.08,0.26,title,size=12,bold=True,color=BG,align=PP_ALIGN.CENTER)
        hdrs=("Function","greedy","samp best")
        col_x=[l+0.05, l+1.60, l+2.65]; col_w=[1.50,1.00,1.00]
        for ri,(row) in enumerate([hdrs]+rows_data):
            y=t+0.38+ri*0.50
            rect(sl,l,y,w,0.47,ACCENT if ri==0 else (PANEL if ri%2 else BG))
            for ci,(cell,cx,cw) in enumerate(zip(row,col_x,col_w)):
                if ri==0: fc=BG
                elif ci>0: fc=ACCENT2 if cell.startswith("+") else RED
                else: fc=LGRAY
                box(sl,cx,y+0.05,cw,0.37,cell,size=11,bold=(ri==0 or ci>0),
                    color=fc,align=PP_ALIGN.CENTER if ci>0 else PP_ALIGN.LEFT)

    tfx_wt=[
        ("array_red.","−0.081","−0.053"),
        ("binary_tree","+0.079","+0.081"),
        ("fft",       "+0.008","+0.009"),
        ("interpreter","+0.307","+0.308"),
        ("kmp_search","+0.077","+0.081"),
        ("poly_eval", "+0.002","+0.036"),
    ]
    gru_wt=[
        ("array_red.","−0.039","−0.011"),
        ("binary_tree","+0.051","+0.058"),
        ("fft",       "+0.016","+0.030"),
        ("interpreter","+0.306","+0.306"),
        ("kmp_search","+0.102","+0.107"),
        ("poly_eval", "+0.118","+0.133"),
    ]
    small_table(0.30,1.20,3.85,"Auto-TFX  Weighted  (mean greedy +6.5%)",ORANGE,tfx_wt)
    small_table(4.30,1.20,3.85,"Auto-GRU  Weighted  (mean greedy +9.2%)",ORANGE,gru_wt)

    # Right: ir_corr plot + explanation
    rect(sl,8.35,1.20,4.75,0.40,RED)
    box(sl,8.39,1.24,4.67,0.30,"IR Reduction vs Runtime Speedup (ir_corr)",
        size=13,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    ir = os.path.join(CKPT,"ir_corr.png")
    if os.path.exists(ir): img(sl,ir,8.35,1.64,4.75)
    bullet_box(sl,8.35,5.38,4.75,1.80,[
        (0,"Top 20 IR-step sequences all from fft"),
        (0,"IR reduction: 65–67%  (near-vertical cluster)"),
        (0,"Runtime speedup: −8% to −35%  — all negative"),
        (0,"Conclusion: IR reduction ≠ runtime improvement"),
        (1,"Stripping fft IR removes vectorization opportunities"),
        (0,"Explains why weighted return underperforms episode"),
    ],size=12,color=LGRAY,rich=True)
content_slide(sl,"Weighted Evaluation & IR-Reduction Correlation",_s19_wt_ir)


# ═══════════════════════════════════════════════════════════════════════════════
# 20 — Architecture × Return Comparison Summary + Pool Result
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
def _s20_summary(sl):
    # Comparison table (all 4 configs)
    box(sl,0.35,1.22,8.10,0.38,"Architecture × Return Summary",
        size=17,bold=True,color=WHITE)
    hdrs=("Architecture","Return","Greedy mean","Train EV")
    data=[
        ("Auto-TFX","Episode",  "+10.4%","0.69"),
        ("Auto-GRU","Episode",  "+11.4%","0.59"),
        ("Auto-TFX","Weighted", "+6.5%", "0.87"),
        ("Auto-GRU","Weighted", "+9.2%", "0.86"),
    ]
    col_x=[0.35,2.70,5.10,7.05]; col_w=[2.30,2.35,1.90,1.35]
    for ri,row in enumerate([hdrs]+data):
        y=1.68+ri*0.55
        episode = ri>0 and row[1]=="Episode"
        weighted= ri>0 and row[1]=="Weighted"
        bg_c = ACCENT if ri==0 else (RGBColor(0x0A,0x22,0x14) if episode else RGBColor(0x22,0x14,0x0A))
        rect(sl,0.35,y,8.10,0.52,bg_c)
        for ci,(cell,cx,cw) in enumerate(zip(row,col_x,col_w)):
            if ri==0: fc=BG
            elif ci==2: fc=ACCENT2 if episode else ORANGE
            elif ci==3: fc=YELLOW
            else: fc=WHITE
            box(sl,cx+0.05,y+0.06,cw-0.10,0.40,cell,size=14,
                bold=(ri==0 or ci==2),color=fc,
                align=PP_ALIGN.CENTER if ci>0 else PP_ALIGN.LEFT)

    bullet_box(sl,0.35,4.50,8.10,2.10,[
        (0,"Episode return outperforms weighted for both architectures"),
        (1,"Weighted EV higher but speedup lower — value head learns IR targets, not speedup"),
        (0,"GRU slightly ahead of TFX in greedy eval despite lower training mean"),
        (1,"Best-checkpoint selection can capture an early peak epoch"),
        (0,"Architecture gap < return-formulation gap — both arches broadly comparable"),
    ],size=14,rich=True)

    # Pool result (right panel)
    rect(sl,8.65,1.22,4.40,0.38,RED)
    box(sl,8.69,1.26,4.32,0.28,"Pool Training (38 functions) — Failed",
        size=13,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    pi = os.path.join(CKPT,"auto-tfx-pool.png")
    if os.path.exists(pi): img(sl,pi,8.65,1.64,4.40)
    bullet_box(sl,8.65,5.38,4.40,1.80,[
        (0,"Trained on all 38 functions simultaneously"),
        (0,"EV = 0.065 at epoch 100 — value head never converged"),
        (0,"Greedy eval: all 38 functions negative (mean −2.50)"),
        (0,"Cause: too few episodes per function to build reward signal"),
        (0,"Key constraint: episodes per function, not total epochs"),
    ],size=12,color=LGRAY,rich=True)
content_slide(sl,"Results Summary & Pool-Training Outcome",_s20_summary)


# ═══════════════════════════════════════════════════════════════════════════════
# 21 — Benchmarking Harness Rigor
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
def _s18(sl):
    box(sl,0.45,1.25,5.9,0.40,"Benchmarking Harness",size=17,bold=True,color=ACCENT)
    bullet_box(sl,0.45,1.72,5.9,2.30,[
        (0,"201 timed iterations per measurement"),
        (1,"10% trimmed mean — removes 10 outliers each tail"),
        (1,"5 warm-up runs before timing begins"),
        (0,"CLOCK_MONOTONIC — monotonic wall-clock, nanosecond resolution"),
        (0,"Serial baselines collected once under stable single-threaded conditions"),
        (0,"Episode collection: parallelised across all CPU cores via Rayon"),
    ],size=14,color=LGRAY,rich=True)
    rect(sl,0.45,4.10,5.9,2.65,PANEL)
    box(sl,0.55,4.16,5.7,0.36,"Parallel Noise Characterisation",size=15,bold=True,color=ACCENT2)
    bullet_box(sl,0.55,4.58,5.7,2.00,[
        (0,"bench-noise: run each binary solo then across 16 Rayon workers simultaneously"),
        (0,"Parallel overhead: small and consistent across all 6 benchmarks"),
        (0,"Noise margin 1.01 on -O3 baseline during training:"),
        (1,"Policy must beat -O3 by > noise floor to receive positive reward"),
        (0,"Prevents noise-driven false positives in the reward signal"),
    ],size=13,color=LGRAY,rich=True)
    box(sl,6.70,1.25,6.2,0.40,"Reproducibility Verification  (diagnose)",size=17,bold=True,color=ACCENT)
    bullet_box(sl,6.70,1.72,6.2,2.85,[
        (0,"Problem: training speedups collected under parallel timing noise"),
        (0,"diagnose re-benchmarks top sequences serially under controlled conditions"),
        (1,"20 independent serial benchmark runs per sequence"),
        (1,"Reports mean, std, and full distribution"),
        (0,"Result: cached training speedups correlate strongly with re-measured values"),
        (0,"Per-sequence distributions are tight → reward signal is reproducible"),
        (0,"Dominant sequences from interpreter and kmp_search — consistent with EDA ceiling"),
    ],size=14,color=LGRAY,rich=True)
    for path,l in [(os.path.join(CKPT,"bench_noise-fft.png"),6.70),
                   (os.path.join(CKPT,"bench_noise-kmp.png"),9.80)]:
        img(sl,path,l,4.72,2.95)
    diag=os.path.join(CKPT,"diagnose.png")
    if os.path.exists(diag): img(sl,diag,6.70,4.72,6.2)
    elif not os.path.exists(os.path.join(CKPT,"bench_noise-fft.png")):
        rect(sl,6.70,4.72,6.2,2.55,RGBColor(0x0A,0x10,0x1C))
        box(sl,6.70,5.55,6.2,0.80,
            "bench_noise-{fft,kmp}.png · diagnose.png\n"
            "[cargo run -- bench-noise  /  cargo run -- plot-diagnose\n"
            " --results checkpoints/auto-tfx-episode-diagnose.json]",
            size=12,color=RGBColor(0x55,0x55,0x77),align=PP_ALIGN.CENTER,italic=True)
content_slide(sl,"Benchmarking Harness Rigor",_s18)


# ═══════════════════════════════════════════════════════════════════════════════
# 19 — Discussion
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
def _s19(sl):
    bullet_box(sl,0.45,1.25,5.90,5.80,[
        (0,"Episode return uses the true runtime signal directly"),
        (1,"Every step in a winning sequence reinforced, even modest ones"),
        (1,"No proxy metric — policy cannot be misled by IR statistics"),
        (1,"Consistently steers toward faster binaries"),
        (0,"Weighted return: proxy misalignment kills runtime gains"),
        (1,"IR reduction rewarded; no gradient signal for actual speedup"),
        (1,"Higher EV (0.87/0.86) but lower speedup (+6.5%/+9.2%) vs episode"),
        (1,"ir_corr confirms: top IR-reducing fft sequences all −8% to −35% runtime"),
        (0,"Transformer vs GRU: both architectures comparable"),
        (1,"GRU: recency bias — free inductive prior, lower training convergence"),
        (1,"TFX: non-local attention — slower training but similar final eval"),
        (1,"Architecture gap < return-formulation gap"),
    ],size=15,rich=True)
    rect(sl,6.50,1.25,6.45,5.80,PANEL)
    box(sl,6.60,1.30,6.25,0.40,"Stop-Token Behaviour",size=17,bold=True,color=YELLOW)
    bullet_box(sl,6.60,1.78,6.25,3.50,[
        (0,"Policy never learns to Stop early — always uses all K=20 steps"),
        (0,"Rational given credit-assignment structure:"),
        (1,"Q(Stop) = current speedup"),
        (1,"Q(any pass) = E[future speedup] — uncertain, possibly higher"),
        (0,"No-op penalty helps at the margin but doesn't fix the asymmetry:"),
        (1,"Cost of stopping: unbounded foregone gain"),
        (1,"Cost of one more pass: at most p = 0.025"),
        (0,"Many top sequences: 5–7 repetitions of dominant pass"),
        (1,"Repeating is safe and occasionally lucky"),
    ],size=13,color=LGRAY,rich=True)
    box(sl,6.60,5.40,6.25,0.36,"Potential Fix:",size=14,bold=True,color=ACCENT2)
    bullet_box(sl,6.60,5.82,6.25,0.90,[
        (0,"Per-step length cost, or Stop bonus ≥ E[marginal return from one more pass]"),
    ],size=13,color=LGRAY,rich=True)
content_slide(sl,"Discussion",_s19)


# ═══════════════════════════════════════════════════════════════════════════════
# 20 — Limitations & Future Work
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
def _s20(sl):
    bullet_box(sl,0.45,1.25,5.90,2.80,[
        (0,"Limitations"),
        (1,"Benchmark timing noise limits per-episode credit precision"),
        (1,"Training set of 6 functions is small"),
        (1,"Pass parameters not tuned (loop-unroll count always default)"),
        (1,"Effective horizon = K; Stop token unused"),
    ],size=16,rich=True)
    bullet_box(sl,0.45,4.30,5.90,2.80,[
        (0,"Key Findings"),
        (1,"Episode return > weighted: IR reduction is a poor runtime proxy"),
        (1,"Pool (38 fn) failed to converge — episodes per function is the constraint"),
        (1,"Architecture gap smaller than return-formulation gap"),
    ],size=16,rich=True)
    bullet_box(sl,6.50,1.25,6.45,5.70,[
        (0,"Future Directions"),
        (1,"Per-step length cost or Stop bonus"),
        (1,"Extend pass menu to cover more LLVM 20 transforms"),
        (1,"Transfer to larger, real-world C/C++ programs"),
        (1,"Beam search or MCTS at inference time"),
        (1,"Combine offline bench-cache with online RL updates"),
        (1,"Pass hyperparameter tuning as part of the action space"),
        (1,"Cross-architecture generalisation (ARM, RISC-V)"),
    ],size=16,rich=True)
content_slide(sl,"Limitations & Future Work",_s20)


# ═══════════════════════════════════════════════════════════════════════════════
# 21 — Conclusion
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK); bg(sl)
rect(sl,0,0,0.10,7.5,ACCENT)
box(sl,0.6,0.5,12.1,0.7,"Conclusion",size=34,bold=True,color=WHITE)
hline(sl,1.35,ACCENT,Pt(1))
bullet_box(sl,0.6,1.52,12.0,5.10,[
    (0,"Episode return beats instruction-weighted for both architectures"),
    (1,"TFX episode: greedy mean +10.4%, 5/6 functions · sample-best all 6 (2–27%)"),
    (1,"GRU episode: greedy mean +11.4% — comparable to TFX despite slower training"),
    (1,"Weighted: +6.5% (TFX) / +9.2% (GRU) — IR-reduction proxy diverges from speedup"),
    (0,"ir_corr confirms IR reduction is not a reliable speedup predictor"),
    (1,"Top IR-reducing fft sequences achieve −8% to −35% runtime — all slower than -O3"),
    (0,"Stop-token analysis: policy always exhausts horizon K=20"),
    (1,"Rational under terminal reward; fix requires per-step length cost or Stop bonus"),
    (0,"Pool-wide training (38 fn) failed to converge — episodes/function is the binding constraint"),
    (0,"Key insight: IR delta features must capture structural position, not global counts"),
    (1,"loop-rotate, licm, sroa shifts invisible to global count vectors; visible to deltas"),
],size=17,color=WHITE,rich=True)
box(sl,0.6,6.6,12.1,0.55,
    "Learning LLVM Pass Sequences via Reinforcement Learning with Autoregressive Policies"
    "  ·  Evan Black  ·  ODU",
    size=13,color=LGRAY,italic=True)


# ── Save ──────────────────────────────────────────────────────────────────────
out = os.path.join(BASE,"presentation.pptx")
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
